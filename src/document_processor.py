import os
import hashlib
import json
from pathlib import Path
import logging
from typing import List, Dict, Tuple, Optional

# Document processing libraries
import PyPDF2
import pdfplumber  # Better for tables and layout
from pptx import Presentation  # PowerPoint
from docx import Document  # Word documents
import openpyxl  # Excel files
from PIL import Image
import pytesseract  # OCR for images

# Text processing
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb
from chromadb.config import Settings

logger = logging.getLogger(__name__)

class AdvancedDocumentProcessor:
    """Enhanced document processor supporting multiple file types with duplicate detection"""
    
    SUPPORTED_EXTENSIONS = {
        'pdf', 'docx', 'pptx', 'xlsx', 'txt', 'md',
        'png', 'jpg', 'jpeg', 'tiff', 'bmp'  # Image files for OCR
    }
    
    def __init__(self, 
                 chroma_client: chromadb.Client,
                 chunk_size: int = 1000,
                 chunk_overlap: int = 200,
                 processed_files_db: str = "processed_files.json"):
        
        self.chroma_client = chroma_client
        self.processed_files_db = processed_files_db
        
        # Text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
        )
        
        # Load processed files registry
        self.processed_files = self._load_processed_files()
    
    def _load_processed_files(self) -> Dict:
        """Load registry of already processed files"""
        try:
            if os.path.exists(self.processed_files_db):
                with open(self.processed_files_db, 'r') as f:
                    return json.load(f)
            return {}
        except Exception as e:
            logger.error(f"Error loading processed files registry: {e}")
            return {}
    
    def _save_processed_files(self):
        """Save processed files registry"""
        try:
            with open(self.processed_files_db, 'w') as f:
                json.dump(self.processed_files, f, indent=2)
        except Exception as e:
            logger.error(f"Error saving processed files registry: {e}")
    
    def _get_file_hash(self, file_path: str) -> str:
        """Generate hash of file content for duplicate detection"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def _is_file_processed(self, file_path: str) -> bool:
        """Check if file was already processed"""
        file_hash = self._get_file_hash(file_path)
        filename = os.path.basename(file_path)
        
        if filename in self.processed_files:
            return self.processed_files[filename].get('hash') == file_hash
        return False
    
    def _mark_file_processed(self, file_path: str, chunks_created: int):
        """Mark file as processed"""
        file_hash = self._get_file_hash(file_path)
        filename = os.path.basename(file_path)
        
        self.processed_files[filename] = {
            'hash': file_hash,
            'processed_at': str(Path(file_path).stat().st_mtime),
            'chunks_created': chunks_created,
            'file_size': os.path.getsize(file_path)
        }
        self._save_processed_files()
    
    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """Enhanced PDF extraction with table and image support"""
        text = ""
        
        logger.info(f"📖 Processing PDF with advanced extraction: {pdf_path}")
        
        try:
            # Method 1: pdfplumber for better table handling
            import pdfplumber
            with pdfplumber.open(pdf_path) as pdf:
                logger.info(f"📄 PDF has {len(pdf.pages)} pages")
                
                for page_num, page in enumerate(pdf.pages, 1):
                    logger.info(f"  📄 Processing page {page_num}...")
                    
                    # Extract regular text
                    page_text = page.extract_text() or ""
                    
                    # Extract tables
                    tables = page.extract_tables()
                    if tables:
                        logger.info(f"    📊 Found {len(tables)} tables on page {page_num}")
                        for table_num, table in enumerate(tables, 1):
                            text += f"\n\n--- TABLE {table_num} (Page {page_num}) ---\n"
                            for row in table:
                                if row:  # Skip empty rows
                                    text += " | ".join([cell or "" for cell in row]) + "\n"
                            text += "--- END TABLE ---\n\n"
                    
                    text += page_text + f"\n[End of Page {page_num}]\n"
                    
        except ImportError:
            logger.warning("pdfplumber not available, falling back to PyPDF2")
            # Fallback to PyPDF2
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        
        except Exception as e:
            logger.error(f"Error in PDF extraction: {e}")
            raise
        
        logger.info(f"📝 Extracted {len(text)} characters from PDF")
        return text
    
    def extract_text_from_docx(self, docx_path: str) -> str:
        """Extract text from Word documents"""
        logger.info(f"📄 Processing Word document: {docx_path}")
        
        try:
            doc = Document(docx_path)
            text = ""
            
            # Extract paragraphs
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Extract tables
            for table_num, table in enumerate(doc.tables, 1):
                text += f"\n\n--- TABLE {table_num} ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    text += row_text + "\n"
                text += "--- END TABLE ---\n\n"
            
            logger.info(f"📝 Extracted {len(text)} characters from Word document")
            return text
            
        except Exception as e:
            logger.error(f"Error extracting from Word document: {e}")
            raise
    
    def extract_text_from_pptx(self, pptx_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        logger.info(f"📊 Processing PowerPoint: {pptx_path}")
        
        try:
            prs = Presentation(pptx_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n\n--- SLIDE {slide_num} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                # Extract notes
                if slide.notes_slide.notes_text_frame.text:
                    text += f"\nSlide Notes: {slide.notes_slide.notes_text_frame.text}\n"
                
                text += f"--- END SLIDE {slide_num} ---\n"
            
            logger.info(f"📝 Extracted {len(text)} characters from PowerPoint")
            return text
            
        except Exception as e:
            logger.error(f"Error extracting from PowerPoint: {e}")
            raise
    
    def extract_text_from_xlsx(self, xlsx_path: str) -> str:
        """Extract text from Excel files"""
        logger.info(f"📊 Processing Excel file: {xlsx_path}")
        
        try:
            wb = openpyxl.load_workbook(xlsx_path)
            text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n\n--- SHEET: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():  # Skip empty rows
                        text += row_text + "\n"
                
                text += f"--- END SHEET: {sheet_name} ---\n"
            
            logger.info(f"📝 Extracted {len(text)} characters from Excel")
            return text
            
        except Exception as e:
            logger.error(f"Error extracting from Excel: {e}")
            raise
    
    def extract_text_from_image(self, image_path: str) -> str:
        """Extract text from images using OCR"""
        logger.info(f"🖼️ Processing image with OCR: {image_path}")
        
        try:
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image)
            
            logger.info(f"📝 Extracted {len(text)} characters from image")
            return f"--- IMAGE: {os.path.basename(image_path)} ---\n{text}\n--- END IMAGE ---\n"
            
        except Exception as e:
            logger.error(f"Error extracting from image: {e}")
            raise
    
    def extract_text_from_file(self, file_path: str) -> str:
        """Universal file text extraction"""
        file_ext = Path(file_path).suffix.lower().lstrip('.')
        
        if file_ext == 'pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_ext == 'docx':
            return self.extract_text_from_docx(file_path) 
        elif file_ext == 'pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_ext == 'xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif file_ext in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
            return self.extract_text_from_image(file_path)
        elif file_ext in ['txt', 'md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    
    def process_file(self, file_path: str, collection_name: str = "documents") -> Dict:
        """Process a single file with duplicate detection"""
        filename = os.path.basename(file_path)
        
        # Check if already processed
        if self._is_file_processed(file_path):
            logger.info(f"⏭️ SKIPPING: '{filename}' already processed")
            return {
                'filename': filename,
                'status': 'skipped',
                'reason': 'already_processed',
                'chunks_created': 0
            }
        
        try:
            logger.info(f"🔄 PROCESSING: '{filename}'")
            
            # Extract text
            text = self.extract_text_from_file(file_path)
            
            # Create chunks
            chunks = self.text_splitter.split_text(text)
            logger.info(f"✂️ Created {len(chunks)} chunks")
            
            # Prepare for ChromaDB
            chunk_objects = []
            for i, chunk in enumerate(chunks):
                chunk_objects.append({
                    'id': f"{filename}_{i}",
                    'text': chunk,
                    'metadata': {
                        'source': filename,
                        'file_type': Path(file_path).suffix.lower(),
                        'chunk_index': i,
                        'total_chunks': len(chunks)
                    }
                })
            
            # Store in ChromaDB
            self._store_chunks_in_chromadb(chunk_objects, collection_name)
            
            # Mark as processed
            self._mark_file_processed(file_path, len(chunks))
            
            return {
                'filename': filename,
                'status': 'processed',
                'chunks_created': len(chunks),
                'characters_extracted': len(text)
            }
            
        except Exception as e:
            logger.error(f"❌ Error processing '{filename}': {e}")
            return {
                'filename': filename,
                'status': 'error',
                'error': str(e),
                'chunks_created': 0
            }
    
    def _store_chunks_in_chromadb(self, chunks: List[Dict], collection_name: str = "documents"):
        """Store chunks in ChromaDB"""
        try:
            collection = self.chroma_client.get_or_create_collection(
                name=collection_name,
                metadata={"description": "RAG document collection"}
            )
            
            documents = [chunk['text'] for chunk in chunks]
            metadatas = [chunk['metadata'] for chunk in chunks]
            ids = [chunk['id'] for chunk in chunks]
            
            collection.add(
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            
            logger.info(f"💾 Stored {len(chunks)} chunks in ChromaDB")
            
        except Exception as e:
            logger.error(f"Error storing chunks: {e}")
            raise
    
    def process_folder(self, folder_path: str, collection_name: str = "documents") -> Dict:
        """Process all supported files in a folder"""
        folder_path = Path(folder_path)
        
        if not folder_path.exists():
            raise ValueError(f"Folder does not exist: {folder_path}")
        
        logger.info(f"📁 PROCESSING FOLDER: {folder_path}")
        
        results = {
            'processed': [],
            'skipped': [],
            'errors': [],
            'total_chunks': 0
        }
        
        # Find all supported files
        supported_files = []
        for ext in self.SUPPORTED_EXTENSIONS:
            supported_files.extend(folder_path.glob(f"*.{ext}"))
            supported_files.extend(folder_path.glob(f"*.{ext.upper()}"))
        
        logger.info(f"📄 Found {len(supported_files)} supported files")
        
        # Process each file
        for file_path in supported_files:
            result = self.process_file(str(file_path), collection_name)
            
            if result['status'] == 'processed':
                results['processed'].append(result)
                results['total_chunks'] += result['chunks_created']
            elif result['status'] == 'skipped':
                results['skipped'].append(result)
            else:
                results['errors'].append(result)
        
        logger.info(f"🎯 FOLDER PROCESSING COMPLETE:")
        logger.info(f"  ✅ Processed: {len(results['processed'])} files")
        logger.info(f"  ⏭️ Skipped: {len(results['skipped'])} files") 
        logger.info(f"  ❌ Errors: {len(results['errors'])} files")
        logger.info(f"  📦 Total chunks: {results['total_chunks']}")
        
        return results 